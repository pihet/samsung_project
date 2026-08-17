import os
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. 차트 이미지 생성 (EDA 슬라이드용)
output_dir = os.path.dirname(os.path.abspath(__file__))
chart1_path = os.path.join(output_dir, "eda_chart1_seasonal.png")
chart2_path = os.path.join(output_dir, "eda_chart2_stationarity.png")

# 한글 폰트 설정
plt.rc('font', family='Malgun Gothic')
plt.rc('axes', unicode_minus=False)

# 차트 1: 7일 주기 계절성 및 Lag 7 자기상관
plt.style.use('default')
plt.rc('font', family='Malgun Gothic')
fig, ax = plt.subplots(figsize=(4.5, 2.2), dpi=200)
days = ['월', '화', '수', '목', '금', '토', '일']
prices = [45.2, 45.0, 44.8, 44.9, 43.6, 43.2, 43.5] # 주말 하락 패턴
bars = ax.bar(days, prices, color=['#CBD5E1', '#CBD5E1', '#CBD5E1', '#CBD5E1', '#EF4444', '#EF4444', '#EF4444'], width=0.55)
ax.set_ylim(41, 47)
ax.set_title('요일별 평균 가격 (금·토·일 주말 특가 하락)', fontsize=9, fontweight='bold', pad=6, color='#1E293B')
ax.tick_params(axis='both', labelsize=8)
for spine in ['top', 'right']:
    ax.spines[spine].set_visible(False)
ax.grid(axis='y', linestyle='--', alpha=0.5)
plt.tight_layout()
plt.savefig(chart1_path, dpi=200, bbox_inches='tight', transparent=True)
plt.close()

# 차트 2: 원시 vs 1차 차분 비교
fig, ax = plt.subplots(1, 2, figsize=(4.5, 2.2), dpi=200)
np.random.seed(42)
t = np.arange(30)
raw_p = 450 - t*1.5 + np.sin(t/2)*10 + np.random.normal(0, 3, 30)
diff_p = np.diff(raw_p)

ax[0].plot(raw_p, color='#115DCE', lw=1.5)
ax[0].set_title('원시 시계열 (비정상)', fontsize=8, fontweight='bold', color='#115DCE')
ax[0].tick_params(labelsize=7)
ax[0].grid(True, linestyle='--', alpha=0.4)

ax[1].plot(diff_p, color='#10B981', lw=1.5)
ax[1].axhline(0, color='gray', linestyle=':', alpha=0.7)
ax[1].set_title('1차 차분 후 (정상성 확보)', fontsize=8, fontweight='bold', color='#10B981')
ax[1].tick_params(labelsize=7)
ax[1].grid(True, linestyle='--', alpha=0.4)
plt.tight_layout()
plt.savefig(chart2_path, dpi=200, bbox_inches='tight', transparent=True)
plt.close()

# 2. PowerPoint 프레젠테이션 생성 (16:9 와이드)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 배경색 (화이트/슬레이트)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
bg.line.fill.background()

# 상단 타이틀 뱃지
badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.65), Inches(0.55))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(17, 93, 206) # #115DCE
badge.line.fill.background()
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "11"
p_b.font.size = Pt(18)
p_b.font.bold = True
p_b.font.color.rgb = RGBColor(255, 255, 255)
p_b.alignment = PP_ALIGN.CENTER

# 타이틀 텍스트
txBox = slide.shapes.add_textbox(Inches(1.55), Inches(0.52), Inches(10.8), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "탐색적 데이터 분석 (EDA) — 시계열 특성 및 모델링 근거 도출"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = RGBColor(15, 23, 42)

# 상단 구분선
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(203, 213, 225)
line.line.fill.background()

# 3개의 카드 컬럼 정의
card_w = Inches(3.7)
card_h = Inches(5.6)
gap = Inches(0.316)
start_x = Inches(0.8)
top_y = Inches(1.4)

card_data = [
    {
        "num": "01",
        "title": "주간 계절성 & 자기상관",
        "tag": "Lag 7 자기상관 r = 0.9057",
        "tag_color": RGBColor(225, 29, 72),
        "img": chart1_path,
        "bullets": [
            "• Lag 7에서 0.9057의 강한 자기상관 관측",
            "• 요일별 분석 결과 주말(금~일) 2~4% 특가 집중",
            "• [모델 반영]: 7일 주기 계절성 가중치 적용"
        ]
    },
    {
        "num": "02",
        "title": "7일 이동 변동성 탐색",
        "tag": "Rolling Volatility 분석",
        "tag_color": RGBColor(139, 92, 246),
        "img": None,
        "custom_content": "volatility_table",
        "bullets": [
            "• 7일 Rolling Std로 단발성 튐 vs 지속 추세 분리",
            "• 정상 변동폭(±1.5%) 초과 급변동 스파이크 감지",
            "• [전처리 반영]: 2.5×IQR 완화 기준 수립 근거"
        ]
    },
    {
        "num": "03",
        "title": "ADF + KPSS 정상성 검정",
        "tag": "1차 차분 후 100% 정상성 확보",
        "tag_color": RGBColor(16, 185, 129),
        "img": chart2_path,
        "bullets": [
            "• 원시 시계열: ADF p=0.8978 ❌ (비정상)",
            "• 1차 차분 시계열: ADF p=0.0000 ✅ (정상)",
            "• [모델 반영]: Auto-ARIMA d=1 차분 채택"
        ]
    }
]

for i, c in enumerate(card_data):
    cx = start_x + i * (card_w + gap)
    
    # 카드 컨테이너 배경
    card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, top_y, card_w, card_h)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card_bg.line.color.rgb = RGBColor(226, 232, 240)
    card_bg.line.width = Pt(1.5)
    
    # 카드 상단 헤더 바
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, top_y, card_w, Inches(0.55))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    header_bar.line.fill.background()
    
    # 카드 번호 및 제목
    tf_h = header_bar.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = f"{c['num']}  {c['title']}"
    p_h.font.size = Pt(13)
    p_h.font.bold = True
    p_h.font.color.rgb = RGBColor(255, 255, 255)
    p_h.alignment = PP_ALIGN.LEFT
    
    # 핵심 지표 뱃지
    tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), top_y + Inches(0.7), card_w - Inches(0.4), Inches(0.42))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    tag_box.line.color.rgb = c['tag_color']
    tag_box.line.width = Pt(1.2)
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = f"★ {c['tag']}"
    p_tag.font.size = Pt(10.5)
    p_tag.font.bold = True
    p_tag.font.color.rgb = c['tag_color']
    p_tag.alignment = PP_ALIGN.CENTER
    
    # 차트 이미지 또는 표 삽입
    if c['img'] and os.path.exists(c['img']):
        slide.shapes.add_picture(c['img'], cx + Inches(0.2), top_y + Inches(1.25), width=card_w - Inches(0.4))
    elif c.get('custom_content') == 'volatility_table':
        tbl_box = slide.shapes.add_table(3, 2, cx + Inches(0.2), top_y + Inches(1.3), card_w - Inches(0.4), Inches(1.6))
        table = tbl_box.table
        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(1.8)
        
        table.cell(0, 0).text = "구분"
        table.cell(0, 1).text = "수치 / 판정"
        table.cell(1, 0).text = "평균 수집 간격"
        table.cell(1, 1).text = "1.9일 (불규칙)"
        table.cell(2, 0).text = "최대 수집 공백"
        table.cell(2, 1).text = "3.9일 (보간 필요)"
        
        for row in table.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9.5)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(51, 65, 85)
                    paragraph.alignment = PP_ALIGN.CENTER
    
    # 하단 텍스트 불릿 포인트
    text_box = slide.shapes.add_textbox(cx + Inches(0.15), top_y + Inches(3.6), card_w - Inches(0.3), Inches(1.8))
    tf_txt = text_box.text_frame
    tf_txt.word_wrap = True
    for b_idx, bullet in enumerate(c['bullets']):
        p_b = tf_txt.add_paragraph() if b_idx > 0 else tf_txt.paragraphs[0]
        p_b.text = bullet
        p_b.font.size = Pt(10.5)
        if "[모델 반영]" in bullet or "[전처리 반영]" in bullet:
            p_b.font.bold = True
            p_b.font.color.rgb = RGBColor(17, 93, 206)
        else:
            p_b.font.color.rgb = RGBColor(71, 85, 105)
        p_b.space_after = Pt(4)

pptx_out = os.path.join(output_dir, "BuyOrWait_EDA_Slide.pptx")
prs.save(pptx_out)
print(f"[SUCCESS] PPTX EDA Slide generated: {pptx_out}")
